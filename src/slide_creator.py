"""
Slide Creator Module
Generates PowerPoint presentations from slides JSON and icons
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path
import json
from typing import Dict, List


class SlideCreator:
    """Create professional PowerPoint presentations"""
    
    def __init__(self, slides_json_path: str = "output/intermediate/slides.json",
                 icons_dir: str = "output/intermediate/images"):
        """
        Initialize slide creator
        
        Args:
            slides_json_path: Path to slides JSON file
            icons_dir: Directory containing icon images
        """
        self.slides_json_path = Path(slides_json_path)
        self.icons_dir = Path(icons_dir)
        
        if not self.slides_json_path.exists():
            raise FileNotFoundError(f"Slides JSON not found: {slides_json_path}")
        
        # Load slides data
        with open(self.slides_json_path, 'r', encoding='utf-8') as f:
            self.data = json.load(f)
        
        # Color scheme (vibrant modern theme)
        self.colors = {
            'primary': RGBColor(41, 98, 255),      # Electric blue
            'secondary': RGBColor(16, 185, 129),   # Emerald green
            'accent': RGBColor(251, 146, 60),      # Warm orange
            'purple': RGBColor(139, 92, 246),      # Purple
            'pink': RGBColor(236, 72, 153),        # Pink
            'bg_light': RGBColor(248, 250, 252),   # Light gray bg
            'text': RGBColor(15, 23, 42),          # Almost black
            'light_text': RGBColor(100, 116, 139), # Slate gray
            'white': RGBColor(255, 255, 255)
        }
    
    def create_presentation(self, output_path: str = "output/slides.pptx") -> str:
        """
        Create complete PowerPoint presentation
        
        Args:
            output_path: Path to save the PPTX file
            
        Returns:
            Path to created presentation
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        print(f"\nCreating PowerPoint presentation...")
        print(f"Output: {output_path}\n")
        
        # Title slide
        print("Creating title slide...")
        self._create_title_slide(prs)
        
        # Content slides
        for i, slide_data in enumerate(self.data['content_slides'], 1):
            print(f"Creating slide {i}: {slide_data['title'][:40]}...")
            self._create_content_slide(prs, slide_data, i)
        
        # Save presentation
        prs.save(str(output_path))
        print(f"\nPresentation created: {output_path}")
        print(f"Total slides: {len(prs.slides)}")
        
        return str(output_path)
    
    def _create_title_slide(self, prs: Presentation):
        """Create vibrant title slide with gradient effect"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        title_data = self.data['title_slide']
        
        # Gradient background using multiple rectangles
        # Top gradient (primary color)
        top_bg = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), prs.slide_width, Inches(4)
        )
        top_bg.fill.solid()
        top_bg.fill.fore_color.rgb = self.colors['primary']
        top_bg.line.fill.background()
        
        # Bottom gradient (secondary color)
        bottom_bg = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(4), prs.slide_width, Inches(3.5)
        )
        bottom_bg.fill.solid()
        bottom_bg.fill.fore_color.rgb = self.colors['secondary']
        bottom_bg.line.fill.background()
        
        # Decorative accent circle (top right)
        circle = slide.shapes.add_shape(
            1,  # Rectangle (will be styled as circle)
            Inches(7.5), Inches(-1), Inches(4), Inches(4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.colors['accent']
        circle.fill.transparency = 0.3
        circle.line.fill.background()
        
        # Decorative accent circle (bottom left)
        circle2 = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(-1.5), Inches(5), Inches(3.5), Inches(3.5)
        )
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = self.colors['purple']
        circle2.fill.transparency = 0.3
        circle2.line.fill.background()
        
        # Icon (centered at top)
        icon_path = self.icons_dir / "title.png"
        if icon_path.exists():
            icon_left = Inches(4)
            icon_top = Inches(1.2)
            icon_width = Inches(2)
            
            slide.shapes.add_picture(
                str(icon_path),
                icon_left, icon_top,
                width=icon_width
            )
        
        # Title with shadow effect
        title_left = Inches(0.5)
        title_top = Inches(3.2)
        title_width = Inches(9)
        title_height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        
        p = title_frame.paragraphs[0]
        p.text = title_data['title']
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle with modern styling
        subtitle_left = Inches(1)
        subtitle_top = Inches(5)
        subtitle_width = Inches(8)
        subtitle_height = Inches(1.2)
        
        subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
        subtitle_frame = subtitle_box.text_frame
        
        p = subtitle_frame.paragraphs[0]
        p.text = title_data['subtitle']
        p.font.size = Pt(28)
        p.font.color.rgb = self.colors['white']
        p.alignment = PP_ALIGN.CENTER
    
    def _create_content_slide(self, prs: Presentation, slide_data: Dict, slide_num: int):
        """Create vibrant content slide with modern design"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Alternate colors for variety
        colors = [self.colors['primary'], self.colors['secondary'],
                 self.colors['purple'], self.colors['pink']]
        slide_color = colors[slide_num % len(colors)]
        
        # Light background
        bg = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors['bg_light']
        bg.line.fill.background()
        
        # Decorative side bar (left accent)
        side_bar = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), Inches(0.4), prs.slide_height
        )
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = slide_color
        side_bar.line.fill.background()
        
        # Icon (compact, left side)
        icon_path = self.icons_dir / f"slide_{slide_num:02d}.png"
        if icon_path.exists():
            icon_left = Inches(1)
            icon_top = Inches(2)
            icon_width = Inches(1.8)
            
            # Icon background circle
            icon_bg = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0.7), Inches(1.7), Inches(2.4), Inches(2.4)
            )
            icon_bg.fill.solid()
            icon_bg.fill.fore_color.rgb = self.colors['white']
            icon_bg.line.color.rgb = slide_color
            icon_bg.line.width = Pt(3)
            
            slide.shapes.add_picture(
                str(icon_path),
                icon_left, icon_top,
                width=icon_width
            )
        
        # Title with colored background
        title_left = Inches(3.5)
        title_top = Inches(0.8)
        title_width = Inches(6)
        title_height = Inches(1.2)
        
        title_bg = slide.shapes.add_shape(
            1,  # Rectangle
            title_left - Inches(0.2), title_top - Inches(0.1),
            title_width + Inches(0.4), title_height + Inches(0.2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = slide_color
        title_bg.line.fill.background()
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = title_frame.paragraphs[0]
        p.text = slide_data['title']
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.colors['white']
        
        # Content area with styled bullets
        content_left = Inches(3.5)
        content_top = Inches(2.5)
        content_width = Inches(6)
        content_height = Inches(4.2)
        
        # Content background (subtle)
        content_bg = slide.shapes.add_shape(
            1,  # Rectangle
            content_left - Inches(0.3), content_top - Inches(0.3),
            content_width + Inches(0.6), content_height + Inches(0.6)
        )
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = self.colors['white']
        content_bg.line.color.rgb = RGBColor(226, 232, 240)
        content_bg.line.width = Pt(2)
        
        content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # Bullets with custom styling
        for i, bullet in enumerate(slide_data['bullets']):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            # Add bullet prefix for visual interest
            bullet_icons = ['-']
            p.text = f"- {bullet}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['text']
            p.level = 0
            p.space_before = Pt(14) if i > 0 else Pt(0)
            p.line_spacing = 1.3
        
        # Decorative corner element (bottom right)
        corner = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(8.5), Inches(6.5), Inches(2), Inches(1.5)
        )
        corner.fill.solid()
        corner.fill.fore_color.rgb = slide_color
        corner.fill.transparency = 0.2
        corner.line.fill.background()
        
        # Add notes (speaker notes)
        if 'speaker_notes' in slide_data and slide_data['speaker_notes']:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data['speaker_notes']


def create_slides(slides_json_path: str = "output/intermediate/slides.json",
                  icons_dir: str = "output/intermediate/images",
                  output_path: str = "output/slides.pptx") -> str:
    """
    Create PowerPoint presentation from slides JSON and icons
    
    Args:
        slides_json_path: Path to slides JSON file
        icons_dir: Directory containing icon images
        output_path: Path to save the PPTX file
        
    Returns:
        Path to created presentation
    """
    creator = SlideCreator(slides_json_path, icons_dir)
    return creator.create_presentation(output_path)


if __name__ == "__main__":
    import sys
    
    try:
        # Create slides
        output_path = create_slides()
        
        print(f"\nSuccess! Open the presentation:")
        print(f"   {output_path}")
        
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
